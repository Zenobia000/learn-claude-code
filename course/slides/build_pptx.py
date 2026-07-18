#!/usr/bin/env python3
"""Assemble the 28-page deck: generated PNG backgrounds + editable Chinese text layer.

Text placement is tuned per page to the actual composition of each generated
background (empty zones located by visual inspection + occupancy analysis).
Output: harness-course-slides-v2.pptx (16:9). Page numbers intentionally omitted.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = Path(__file__).parent
GEN = HERE / "generated"
OUT = HERE / "harness-course-slides-v2.pptx"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

BG = RGBColor(0x17, 0x1C, 0x28)
FG = RGBColor(0xE6, 0xE9, 0xF0)
AMBER = RGBColor(0xFF, 0xB4, 0x54)
CYAN = RGBColor(0x5C, 0xCF, 0xE6)
MUTE = RGBColor(0x4A, 0x55, 0x68)

ZH = "Microsoft JhengHei"
MONO = "Consolas"

COPYRIGHT = "© 2026 桑尼資料科學"
LEGAL = (
    "本教材（含投影片、影片、講義、範例程式與其編排設計）由桑尼資料科學"
    "（Sunny Data Science）研發，受著作權法及相關智慧財產權法律保護。"
    "© 2026 桑尼資料科學，版權所有。\n\n"
    "未經著作權人事前書面同意，不得以任何形式重製、改作、散布、公開傳輸、"
    "公開播送或作商業使用。本教材僅授權課程學員為個人學習目的使用。\n\n"
    "教材中引用之開源專案（learn-claude-code）依其原授權條款（見該 repo "
    "LICENSE）使用，其著作權歸原作者所有。"
)

S01_CODE = '''def agent_loop(query):
    messages = [{"role": "user", "content": query}]
    while True:
        response = client.messages.create(
            model=MODEL, system=SYSTEM,
            messages=messages, tools=TOOLS, max_tokens=8000)
        messages.append(
            {"role": "assistant", "content": response.content})
        if response.stop_reason != "tool_use":
            return
        results = [run_tool(b) for b in response.content
                   if b.type == "tool_use"]
        messages.append({"role": "user", "content": results})'''

# Per-page geometry keys (inches):
#   kx,ky        kicker position          tx,ty,tw,tsz  title box / size
#   bx,by,bw,bsz body box / font size     fx,fy         foot position
# Placement rationale: fitted to each background's empty zones.
SLIDES = [
    dict(page="p01", layout="cover",
         kicker="桑尼資料科學 × Agent 工程系列",
         title="Coding Agent 與\nHarness 工程",
         sub="從打 API 的人，變成設計 Agent Loop 的人",
         foot="教材基於開源專案 learn-claude-code"),
    # p02: art (user->cloud) upper-right band; text lower-left half
    dict(page="p02", layout="std", kicker="第一幕・觀念",
         title="你們熟悉的世界",
         body=["組 prompt → 打 API → 拿回應，一來一回就結束",
               "模型只「回答」，不「行動」：不會讀你的檔案、跑你的測試",
               "Fine-tuning 改變的是模型會說什麼，沒有改變它能做什麼"],
         bx=0.6, by=3.0, bw=6.6, bsz=19),
    # p03: "?" art center at y4-6; question goes top
    dict(page="p03", layout="pause",
         q="你呼叫過的 LLM API，\n為什麼不會自己改你的程式碼？",
         hint="想 30 秒 -- 它缺的不是智力，是什麼？"),
    # p04: car occupies y4.0-6.2; text upper-left
    dict(page="p04", layout="std", kicker="第一幕・觀念",
         title="Agent 產品 = 模型 + Harness",
         body=["Agency（感知、推理、行動）來自模型訓練 -- 只需要這一句結論",
               "模型 = 駕駛者：負責決策與推理",
               "Harness = 載具：工具、知識、權限，讓決策落地成行動",
               "缺一不可 -- 這門課教你造載具"],
         bx=0.6, by=1.75, bw=7.2, bsz=18),
    # p05: hub diagram left; body moves to right column
    dict(page="p05", layout="std", kicker="第一幕・觀念",
         title="Harness 五要素",
         body=["Tools 工具：檔案讀寫 / Shell / 網路",
               "Knowledge 知識：產品文件 / API 規範",
               "Observation 觀察：git diff / 錯誤日誌",
               "Action 行動：CLI 命令 / API 呼叫",
               "Permissions 權限：沙箱隔離 / 信任邊界"],
         bx=6.9, by=2.15, bw=6.1, bsz=18),
    # p06: split art fills center band; body drops to bottom strip
    dict(page="p06", layout="std", kicker="第一幕・觀念",
         title="反例：Prompt Chain 不是 Agent",
         body=["拖拽節點圖、if-else 串 API -- 行為是「編」出來的，不是「學」出來的",
               "每個分支都是人先想好的，遇到沒想到的情況就斷",
               "Agent 的控制流由模型當下決定，harness 只提供行動空間"],
         bx=0.6, by=5.5, bw=12.2, bsz=16),
    # p07: brackets art at y4.4-5.4; question goes top
    dict(page="p07", layout="pause",
         q="那最小可行的 Agent，\n程式碼長什麼樣？",
         hint="先猜一個行數。"),
    # p08: ONE LOOP art on left half; text on right half
    dict(page="p08", layout="section",
         title="一個迴圈，30 行。",
         sub="One loop & Bash is all you need\nagents/s01_agent_loop.py"),
    # p09: flow diagram upper-right; body lower-left
    dict(page="p09", layout="std", kicker="第二幕・作法",
         title="Agent Loop：唯一的控制流",
         body=["messages 累積完整對話，連同工具定義發給模型",
               "stop_reason == \"tool_use\" → 執行工具，結果塞回 messages",
               "stop_reason != \"tool_use\" → 模型決定停，迴圈結束",
               "是模型決定何時停 -- 不是你的 if-else"],
         bx=0.6, by=3.1, bw=6.4, bsz=17),
    # p10: editor window mockup; title + code go inside the window
    dict(page="p10", layout="code", kicker="第二幕・作法",
         title="s01：30 行的完整 Agent", code=S01_CODE,
         foot="之後 11 個機制全部疊在這個迴圈上 -- 迴圈本身不再改動"),
    # p11: terminal window mockup; text goes inside the window
    dict(page="p11", layout="std", kicker="現場 Demo",
         title="Demo：s01 跑通",
         body=["$ python agents/s01_agent_loop.py",
               "任務：Create hello.py that prints \"Hello, World!\"",
               "觀察：模型自己呼叫 bash → 看結果 → 決定停止"],
         kx=3.3, ky=1.25, tx=3.3, ty=1.7, tw=7.0, tsz=30,
         bx=3.5, by=3.0, bw=6.8, bsz=18),
    # p12: staircase fills center; body bottom-right empty zone
    dict(page="p12", layout="std", kicker="第二幕・作法",
         title="機制地圖：12 個 Session，迴圈不變",
         body=["核心骨架 s01–s03：迴圈、工具、計劃",
               "進階機制 s04–s07：子代理、技能、壓縮、任務",
               "多代理協作 s08–s12：背景、團隊、協議、自主、隔離",
               "每個 session 只加一個機制 -- diff 前後兩支檔案就是最好的教材"],
         bx=5.2, by=5.15, bw=7.6, bsz=16),
    # p13: two panels center; body bottom strip
    dict(page="p13", layout="std", kicker="第二幕・作法",
         title="骨架層：工具與計劃（s02–s03）",
         body=["s02 加工具 = 加一個 handler：read / write / edit_file",
               "工具越專用，模型選擇越準 -- loop 完全不用改",
               "s03 先寫計劃再動手：沒有計劃的 agent 會在長任務中漂移 -- 這就是 TodoWrite"],
         bx=0.6, by=5.6, bw=12.2, bsz=16),
    # p14: gauge art center at y3.9-6.6; question goes top
    dict(page="p14", layout="pause",
         q="任務越做越長，\ncontext 會發生什麼事？",
         hint="提示：messages 只進不出。"),
    # p15: three mini-diagrams top-left band; title top-right, body below art
    dict(page="p15", layout="std", kicker="第二幕・作法",
         title="Context 工程（s04–s06）",
         kx=6.9, ky=0.65, tx=6.9, ty=1.05, tw=6.3, tsz=30,
         body=["s04 Subagent：大任務拆給子代理，各用乾淨 context，只回報結論",
               "s05 Skill Loading：知識需要時才載入，不是一開始全塞",
               "s06 Compact：context 快滿時，把歷史壓縮成摘要騰空間"],
         bx=0.9, by=3.4, bw=11.5, bsz=19),
    # p16: three empty cards drawn in art; one mapping per card
    dict(page="p16", layout="cards",
         title="對應到你每天用的 Claude Code",
         cards=[("s04 Subagent", "Task tool（子代理）"),
                ("s05 Skill Loading", "Skills / CLAUDE.md"),
                ("s06 Compact", "/compact 指令")],
         foot="教材裡 200 行內的機制，就是產品裡同名功能的原理"),
    # p17: kanban board left 2/3; body right column
    dict(page="p17", layout="std", kicker="第二幕・作法",
         title="任務系統與背景執行（s07–s08）",
         body=["s07 Task System：檔案式任務圖 -- 可排序、有依賴、可持久化，agent 重啟也不會忘",
               "s08 Background Tasks：慢操作丟背景不阻塞，完成時以通知回到迴圈"],
         bx=0.8, by=5.55, bw=11.9, bsz=16),
    # p18: robots at center y4.9-6.3; question goes top
    dict(page="p18", layout="pause", hy=2.5,
         q="一個 Agent 不夠用的時候呢？",
         hint="多個 agent 一起工作，會遇到哪三個問題？"),
    # p19: lead robot lower-left, workers right column; body upper-left
    dict(page="p19", layout="std", kicker="第二幕・作法 / Multi-Agent",
         title="Multi-Agent I：委派（s09）",
         body=["任務太大 → 委派給隊友，每個隊友是獨立的 agent loop",
               "各自擁有乾淨 context -- 不共享記憶，只交換訊息",
               "非同步信箱：送出任務後不等待，繼續自己的工作"],
         bx=0.6, by=1.85, bw=7.0, bsz=17),
    # p20: two busy scenes; title in top strip, body over lower board
    dict(page="p20", layout="std", kicker="第二幕・作法 / Multi-Agent",
         title="Multi-Agent II：協議與自主（s10–s11）",
         kx=0.6, ky=0.08, tx=0.6, ty=0.42, tw=10.0, tsz=26,
         body=["s10 Team Protocols：共同的訊息格式與回報規則 -- 否則各說各話",
               "s11 Autonomous Agents：不派工，隊友自己掃看板、認領、上工",
               "從「指揮」到「協作」：中心化派工 → 去中心化認領"],
         bx=3.2, by=5.78, bw=9.9, bsz=15),
    # p21: empty banner box drawn at top; title inside it, body above merge line
    dict(page="p21", layout="std", kicker="第二幕・作法 / Multi-Agent",
         title="Multi-Agent III：隔離（s12）",
         kx=0.55, ky=0.14, tx=0.55, ty=0.48, tw=12.4, tsz=28,
         body=["兩個 agent 同時改同一個檔案 = 打架；git worktree：每人一份獨立目錄",
               "各做各的，完成後由 git merge 收斂 -- 這就是 worktree isolation"],
         bx=0.7, by=4.55, bw=12.0, bsz=16),
    # p22: layer stack center-right; body in empty left column
    dict(page="p22", layout="std", kicker="第二幕・作法",
         title="s_full.py：疊完 12 個機制的全貌",
         tx=0.6, ty=0.9, tw=8.5, tsz=32,
         body=["740 行 = 一個 Claude Code 雛形",
               "核心迴圈仍然是 s01 那 30 行",
               "其餘 700 行全是 harness：工具、技能、任務、團隊、隔離"],
         bx=0.6, by=2.5, bw=3.7, bsz=16),
    # p23: three-pane terminal upper; body below window
    dict(page="p23", layout="std", kicker="現場 Demo",
         title="Demo：s_full.py 委派任務",
         body=["$ python agents/s_full.py",
               "任務：一個需要拆解、委派子代理的多步驟任務",
               "觀察：計劃 → 拆任務 → 委派 → 回收結果"],
         bx=3.0, by=5.3, bw=7.6, bsz=16),
    # p24: timeline art in upper third; big statement below the amber line
    dict(page="p24", layout="insight", kicker="第三幕・目的",
         big="12 個機制疊完，\nAgent Loop 一行都沒改。",
         sub="智慧在模型裡，工程在 harness 裡",
         gy=3.5, sy=5.85),
    # p25: chips + chassis fill center; body in bottom caption strip
    dict(page="p25", layout="std", kicker="第三幕・目的",
         title="換模型，不換 Harness",
         body=["改 .env 一行：Claude ↔ MiniMax ↔ GLM ↔ Kimi ↔ DeepSeek",
               "同一套 harness、同一個迴圈，直接開工 -- 載具是你的資產，駕駛者可以換"],
         bx=1.0, by=6.05, bw=11.6, bsz=15),
    # p26: road on left; body in empty right third
    dict(page="p26", layout="std", kicker="第三幕・目的",
         title="課後實作與自學地圖",
         kx=6.6, ky=0.65, tx=6.6, ty=1.05, tw=6.4, tsz=34,
         body=["作業三梯度：生成專案文件 → 補測試 → 實作新功能（用 s_full.py）",
               "影片 14 集：每集一個機制的深度 walkthrough",
               "最小路徑：s01 → s04 → s07，三個檔案掌握 80% 精髓"],
         bx=6.6, by=2.4, bw=6.3, bsz=17),
    # p27: brain + car fill middle band; big statement in lower third
    dict(page="p27", layout="insight",
         big="Agency 是學出來的，Harness 是工程出來的。",
         sub="模型是駕駛者，載具由你打造",
         gy=5.35, gsz=32, sy=6.5),
    dict(page="p28", layout="legal", title="Q&A", legal=LEGAL),
]


def set_alpha(shape, alpha_pct):
    """Set fill transparency (0=opaque, 100=invisible) via raw XML."""
    alpha = str(int((100 - alpha_pct) * 1000))
    srgb = shape.fill.fore_color._xFill.find(qn('a:srgbClr'))
    el = srgb.makeelement(qn('a:alpha'), {'val': alpha})
    srgb.append(el)


def add_text(slide, x, y, w, h, runs, size, color=FG, bold=False, font=ZH,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0,
             backdrop=False):
    if backdrop:
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = BG
        set_alpha(rect, 25)
        rect.line.fill.background()
        rect.shadow.inherit = False
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = runs if isinstance(runs, list) else [runs]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(size * 0.45)
        p.text = line
        for r in p.runs:
            r.font.name = font
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font._rPr.set('spc', '20')
            ea = r.font._rPr.makeelement(qn('a:ea'), {'typeface': font})
            r.font._rPr.append(ea)
    return box


def add_fixtures(slide):
    """Bottom-left Sunny Data Science banner + bottom-right copyright tag."""
    strip = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.25), Inches(7.02),
        Inches(2.75), Inches(0.34))
    strip.fill.solid()
    strip.fill.fore_color.rgb = BG
    set_alpha(strip, 15)
    strip.line.color.rgb = MUTE
    strip.line.width = Pt(0.75)
    strip.shadow.inherit = False
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.25), Inches(7.02),
        Inches(0.06), Inches(0.34))
    accent.fill.solid()
    accent.fill.fore_color.rgb = AMBER
    accent.line.fill.background()
    accent.shadow.inherit = False
    add_text(slide, Inches(0.42), Inches(7.02), Inches(2.6), Inches(0.34),
             "桑尼資料科學  SUNNY DATA SCIENCE", 10, color=FG, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(10.0), Inches(7.06), Inches(3.1), Inches(0.3),
             COPYRIGHT, 9, color=MUTE, align=PP_ALIGN.RIGHT,
             anchor=MSO_ANCHOR.MIDDLE)


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    blank = prs.slide_layouts[6]

    for spec in SLIDES:
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BG
        pngs = sorted(GEN.glob(f"{spec['page']}_*.png"))
        if pngs:
            slide.shapes.add_picture(str(pngs[-1]), 0, 0, SLIDE_W, SLIDE_H)
        lay = spec["layout"]

        if lay == "cover":
            add_text(slide, Inches(0.7), Inches(1.1), Inches(7.5), Inches(0.5),
                     spec["kicker"], 16, color=AMBER, bold=True)
            add_text(slide, Inches(0.7), Inches(1.8), Inches(8.0), Inches(2.8),
                     spec["title"].split("\n"), 52, bold=True, spacing=1.05)
            add_text(slide, Inches(0.7), Inches(4.6), Inches(8.0), Inches(0.7),
                     spec["sub"], 22, color=CYAN)
            add_text(slide, Inches(0.7), Inches(5.5), Inches(8.0), Inches(0.5),
                     spec["foot"], 13, color=MUTE)
        elif lay == "pause":
            # all four pause backgrounds keep their art in the lower half
            add_text(slide, Inches(1.2), Inches(0.95), Inches(10.9), Inches(1.9),
                     spec["q"].split("\n"), 36, bold=True,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            add_text(slide, Inches(1.2), Inches(spec.get("hy", 3.0)),
                     Inches(10.9), Inches(0.6), spec["hint"], 17, color=AMBER,
                     align=PP_ALIGN.CENTER)
        elif lay == "section":
            # ONE LOOP art occupies the left half; text on the right half
            add_text(slide, Inches(7.0), Inches(2.5), Inches(6.0), Inches(1.1),
                     spec["title"], 44, bold=True, align=PP_ALIGN.CENTER,
                     backdrop=True)
            add_text(slide, Inches(7.0), Inches(3.8), Inches(6.0), Inches(1.0),
                     spec["sub"].split("\n"), 15, color=MUTE,
                     align=PP_ALIGN.CENTER)
        elif lay == "code":
            # title + code live inside the editor-window mockup (x2.6-10.6)
            add_text(slide, Inches(3.2), Inches(0.72), Inches(7.2), Inches(0.35),
                     spec["kicker"], 13, color=AMBER, bold=True)
            add_text(slide, Inches(3.2), Inches(1.1), Inches(7.2), Inches(0.7),
                     spec["title"], 27, bold=True)
            add_text(slide, Inches(3.3), Inches(2.0), Inches(7.1), Inches(4.2),
                     spec["code"].split("\n"), 12.5, color=FG, font=MONO,
                     spacing=1.0)
            add_text(slide, Inches(3.2), Inches(6.42), Inches(7.2), Inches(0.4),
                     spec["foot"], 13, color=AMBER)
        elif lay == "cards":
            # three empty cards drawn in the art; fill one mapping per card
            add_text(slide, Inches(2.9), Inches(0.02), Inches(7.5), Inches(0.45),
                     spec["title"], 22, bold=True, align=PP_ALIGN.CENTER,
                     backdrop=True)
            zones = [(0.78, 3.55), (4.88, 3.65), (9.08, 3.75)]
            for (cx, cw), (mech, prod) in zip(zones, spec["cards"]):
                add_text(slide, Inches(cx), Inches(2.5), Inches(cw), Inches(0.6),
                         mech, 20, bold=True, align=PP_ALIGN.CENTER)
                add_text(slide, Inches(cx), Inches(3.3), Inches(cw), Inches(0.5),
                         "↓", 20, color=MUTE, align=PP_ALIGN.CENTER)
                add_text(slide, Inches(cx), Inches(4.0), Inches(cw), Inches(0.9),
                         prod, 19, color=CYAN, align=PP_ALIGN.CENTER)
            add_text(slide, Inches(0.6), Inches(6.62), Inches(12.2), Inches(0.35),
                     spec["foot"], 14, color=AMBER, align=PP_ALIGN.CENTER)
        elif lay == "insight":
            if spec.get("kicker"):
                add_text(slide, Inches(0.6), Inches(0.45), Inches(8.0),
                         Inches(0.4), spec["kicker"], 14, color=AMBER, bold=True)
            gy, gsz = spec.get("gy", 2.4), spec.get("gsz", 40)
            add_text(slide, Inches(1.0), Inches(gy), Inches(11.3), Inches(1.9),
                     spec["big"].split("\n"), gsz, bold=True,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                     backdrop=True)
            add_text(slide, Inches(1.0), Inches(spec.get("sy", 5.1)),
                     Inches(11.3), Inches(0.6), spec["sub"], 18, color=CYAN,
                     align=PP_ALIGN.CENTER)
        elif lay == "legal":
            add_text(slide, Inches(0.7), Inches(0.7), Inches(6.0), Inches(1.0),
                     spec["title"], 44, bold=True)
            add_text(slide, Inches(0.7), Inches(2.1), Inches(11.0), Inches(4.1),
                     spec["legal"].split("\n\n"), 14, color=MUTE, spacing=1.15,
                     backdrop=True)
        else:  # std with per-page geometry
            kx, ky = spec.get("kx", 0.6), spec.get("ky", 0.45)
            tx, ty = spec.get("tx", 0.6), spec.get("ty", 0.9)
            tw, tsz = spec.get("tw", 11.5), spec.get("tsz", 34)
            add_text(slide, Inches(kx), Inches(ky), Inches(9.0), Inches(0.38),
                     spec["kicker"], 14, color=AMBER, bold=True)
            add_text(slide, Inches(tx), Inches(ty), Inches(tw), Inches(0.95),
                     spec["title"], tsz, bold=True)
            bx, by = spec.get("bx", 0.6), spec.get("by", 2.15)
            bw, bsz = spec.get("bw", 7.4), spec.get("bsz", 20)
            body_h = Inches((bsz / 20 * 0.62) * len(spec["body"]) + 0.35)
            add_text(slide, Inches(bx), Inches(by), Inches(bw), body_h,
                     spec["body"], bsz, spacing=1.05, backdrop=True)
        add_fixtures(slide)

    prs.save(OUT)
    print(f"saved {OUT} with {len(SLIDES)} slides")


if __name__ == "__main__":
    build()
